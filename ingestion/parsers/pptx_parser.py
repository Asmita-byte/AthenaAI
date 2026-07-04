from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from backend.config import get_settings
from backend.core.logging import get_logger
from ingestion.parsers.base_parser import (
    BaseParser,
    ParsedChunk,
    ParsedDocument,
    ParsedImage,
    ParsedTable,
)

settings = get_settings()
logger = get_logger(__name__)


class PPTXParser(BaseParser):

    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    def parse(self, file_path: Path) -> ParsedDocument:
        logger.info("Parsing PPTX", filename=file_path.name)

        prs = Presentation(str(file_path))

        doc = ParsedDocument(
            filename=file_path.name,
            file_extension=".pptx",
            total_pages=len(prs.slides),
            metadata=self._extract_metadata(prs, file_path),
        )

        doc.text_chunks = self._extract_text(prs, file_path)

        if self.extract_tables:
            doc.table_chunks = self._extract_tables(prs, file_path)

        if self.extract_images:
            doc.image_chunks = self._extract_images(prs, file_path)

        logger.info(
            "PPTX parsing complete",
            filename=file_path.name,
            summary=doc.summary(),
        )

        return doc

    def _extract_metadata(self, prs: Presentation, file_path: Path) -> dict:
        props = prs.core_properties
        return {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "total_pages": len(prs.slides),
            "created": str(props.created or ""),
            "modified": str(props.modified or ""),
        }

    def _extract_text(self, prs: Presentation, file_path: Path) -> list[ParsedChunk]:
        chunks = []
        chunk_index = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            slide_title = ""

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                if shape.shape_type == 13:
                    continue

                is_title = False
                try:
                    if shape.placeholder_format is not None and shape.placeholder_format.idx == 0:
                        is_title = True
                except ValueError:
                    is_title = False

                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if is_title and not slide_title:
                        slide_title = text
                    else:
                        slide_texts.append(text)

            if slide_title or slide_texts:
                content = f"Slide {slide_num}: {slide_title}\n" + "\n".join(slide_texts)
                chunks.append(
                    ParsedChunk(
                        content=content.strip(),
                        chunk_type="text",
                        page_number=slide_num,
                        chunk_index=chunk_index,
                        source_filename=file_path.name,
                        section_title=slide_title,
                        metadata={"slide": slide_num, "title": slide_title},
                    )
                )
                chunk_index += 1

        return chunks

    def _extract_tables(self, prs: Presentation, file_path: Path) -> list[ParsedChunk]:
        chunks = []
        chunk_index = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_table:
                    continue

                table = shape.table
                rows_data = []

                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    rows_data.append(row_data)

                if not rows_data:
                    continue

                headers = rows_data[0]
                rows = rows_data[1:]

                parsed_table = ParsedTable(
                    table_index=chunk_index,
                    page_number=slide_num,
                    headers=headers,
                    rows=rows,
                    raw_text="",
                )

                chunks.append(
                    ParsedChunk(
                        content=parsed_table.to_text(),
                        chunk_type="table",
                        page_number=slide_num,
                        chunk_index=chunk_index,
                        source_filename=file_path.name,
                        table=parsed_table,
                        metadata={"slide": slide_num},
                    )
                )
                chunk_index += 1

        return chunks

    def _extract_images(self, prs: Presentation, file_path: Path) -> list[ParsedChunk]:
        chunks = []
        chunk_index = 0
        images_dir = settings.storage_images_dir

        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.shape_type != 13:
                    continue

                try:
                    image = shape.image
                    image_ext = image.ext
                    image_bytes = image.blob

                    image_filename = (
                        f"{file_path.stem}_slide{slide_num}_img{chunk_index}.{image_ext}"
                    )
                    image_path = images_dir / image_filename

                    with open(image_path, "wb") as f:
                        f.write(image_bytes)

                    parsed_image = ParsedImage(
                        image_index=chunk_index,
                        page_number=slide_num,
                        image_path=str(image_path),
                        image_type="image",
                    )

                    chunks.append(
                        ParsedChunk(
                            content=f"Image extracted from slide {slide_num}.",
                            chunk_type="image",
                            page_number=slide_num,
                            chunk_index=chunk_index,
                            source_filename=file_path.name,
                            image=parsed_image,
                            metadata={"slide": slide_num},
                        )
                    )
                    chunk_index += 1

                except Exception as e:
                    logger.warning(
                        "Failed to extract image",
                        slide=slide_num,
                        error=str(e),
                    )

        return chunks
